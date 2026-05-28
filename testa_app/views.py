import os
import sys
from django.db.models import Count
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.auth import login, authenticate
from django.views.decorators.http import require_http_methods
from django.contrib.auth.forms import AuthenticationForm
from django.contrib import messages
from django.http import JsonResponse


from .forms import UserRegisterForm, PDFUploadForm
from .models import PDFDocument, QuestionAnswer, UserAnalytics, DailyActivity, Bookmark, Quiz, Flashcard, TopicMastery
from .utils import (
    get_file_text, get_text_chunks, get_vector_store, 
    load_vector_store, get_conversational_chain
)
from dotenv import load_dotenv
from datetime import datetime, date
from django.utils import timezone
from .bytez_client import get_bytez_client

load_dotenv()


def _safe_log(msg, *args):
    """Log to console without raising UnicodeEncodeError on Windows (charmap)."""
    try:
        safe = (str(msg).encode('ascii', errors='replace').decode('ascii'),) + tuple(
            str(a).encode('ascii', errors='replace').decode('ascii') for a in args
        )
        print(*safe)
    except Exception:
        print("[log omitted]")



import requests
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from .models import PDFDocument, QuestionAnswer




import requests
from bs4 import BeautifulSoup
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from .models import PDFDocument, QuestionAnswer
import requests
from bs4 import BeautifulSoup

def web_scrape_search(query):
    search_url = f"https://www.google.com/search?q={query.replace(' ', '+')}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3"
    }
    
    response = requests.get(search_url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    snippet = None
    try:
        snippet = soup.find("span", {"class": "aCOpRe"})
        if snippet:
            return snippet.text
        else:
            snippet = soup.find("div", {"class": "BNeawe s3v9rd AP7Wnd"})
            if snippet:
                return snippet.text
    except AttributeError:
        pass

    return "Sorry, I couldn't find an answer to your question."




@login_required
def question_answer(request):
    start_time = timezone.now()
    
    if request.method == 'POST':
        try:
            user_question = request.POST.get('question')
            if not user_question:
                return JsonResponse({"error": "Question is required"}, status=400)
            
            _safe_log("Received question: " + (user_question[:200] if user_question else ""))

            # Load vector store and search (uses absolute path; includes all uploaded docs)
            answer = None
            context = ""
            try:
                vector_store = load_vector_store()
                if vector_store:
                    docs = vector_store.similarity_search(user_question, k=3)
                    context = "\n\n".join([doc.page_content for doc in docs]).strip()
                    if context:
                        chain = get_conversational_chain()
                        answer = chain(context, user_question)
            except Exception as e:
                _safe_log(f"Error with vector store: {e}")
            
            # If no answer from document context, use Bytez with clear instructions when no context
            if not answer:
                try:
                    client = get_bytez_client()
                    if not context:
                        # No uploaded docs or no relevant passages — avoid "I don't have access to files"
                        answer = client.answer_question(
                            user_question,
                            context="No document context was found. Reply helpfully: if they asked about an uploaded document, say you couldn't find relevant content in their uploaded documents and suggest they upload the file first or rephrase; otherwise answer from your general knowledge. Do not say you cannot access files or external documents."
                        )
                    else:
                        answer = client.answer_question(user_question, context=context)
                except ValueError as e2:
                    # API key missing
                    _safe_log(f"API Key Error: {e2}")
                    answer = "API configuration error. Please contact support."
                except Exception as e2:
                    _safe_log(f"Error with Bytez API: {e2}")
                    import traceback
                    try:
                        traceback.print_exc()
                    except UnicodeEncodeError:
                        pass
                    # Try web scrape as fallback
                    try:
                        answer = web_scrape_search(user_question)
                        if not answer or "Sorry" in answer:
                            # Provide helpful error message
                            error_msg = str(e2)
                            if "timeout" in error_msg.lower():
                                answer = "⏱️ The AI service is taking too long to respond. This might be due to high traffic. Please try again in a moment."
                            elif "connection" in error_msg.lower() or "fetch failed" in error_msg.lower():
                                answer = "🔌 Unable to connect to the AI service. Please check your internet connection and try again."
                            elif "401" in error_msg or "403" in error_msg:
                                answer = "🔐 Authentication error with the AI service. Please contact support."
                            elif "500" in error_msg or "server error" in error_msg.lower() or "inference failed" in error_msg.lower():
                                answer = "⚠️ The AI service is temporarily unavailable. The server is experiencing issues. Please try again in a few moments, or try rephrasing your question."
                            else:
                                answer = "⚠️ I'm having trouble connecting to the AI service right now. This is usually temporary. Please try again in a moment, or try rephrasing your question."
                    except Exception as e3:
                        _safe_log(f"Web scrape also failed: {e3}")
                        answer = "⚠️ I'm unable to process your question right now. The AI service and fallback options are temporarily unavailable. Please try again in a few moments."
            
            if answer == "The answer is not available in the context.":
                answer = web_scrape_search(user_question)
                if not answer:
                    answer = "Sorry, I couldn't find an answer to your question."
            
            if not answer or not str(answer).strip():
                answer = "Sorry, I couldn't generate a response. Please try again or rephrase your question."
            
            # Calculate response time
            response_time = (timezone.now() - start_time).total_seconds()
            _safe_log("Response: " + (str(answer)[:100] if answer else "Empty response"))
            
            # Create or update Q&A entry (avoid UNIQUE constraint errors on (user, question))
            from django.db import IntegrityError
            try:
                qa_entry, created = QuestionAnswer.objects.update_or_create(
                    user=request.user,
                    question=user_question,
                    defaults={
                        'answer': str(answer),
                        'response_time': response_time,
                    },
                )
            except IntegrityError:
                # As a fallback, fetch existing and update in-place
                qa_entry = QuestionAnswer.objects.get(user=request.user, question=user_question)
                qa_entry.answer = str(answer)
                qa_entry.response_time = response_time
                qa_entry.save()
                created = False
            
            # Update analytics (only count new Q&A rows toward daily questions / streak)
            try:
                _update_daily_activity(request.user, qa_is_new=created)
                _update_user_analytics(request.user, qa_was_created=created)
            except Exception as e:
                _safe_log(f"Error updating analytics: {e}")
            
            return JsonResponse({"response": answer})
            
        except Exception as e:
            _safe_log(f"Unexpected error in question_answer: {e}")
            import traceback
            try:
                traceback.print_exc()
            except UnicodeEncodeError:
                pass
            err_msg = str(e)
            if "charmap" in err_msg or "codec" in err_msg or isinstance(e, (UnicodeEncodeError, UnicodeDecodeError)):
                err_msg = "A character encoding issue occurred. Try rephrasing your question or uploading a document without special characters."
            return JsonResponse({
                "error": f"An error occurred: {err_msg}"
            }, status=500)

    qa_list = QuestionAnswer.objects.filter(user=request.user).order_by('-created_at')[:10]
    from django.urls import reverse
    try:
        upload_ajax_url = reverse('upload_document_ajax')
    except Exception:
        upload_ajax_url = '/question_answer/upload/'
    return render(request, 'testa_app/question_answer.html', {
        'qa_list': qa_list,
        'username': request.user.username,
        'upload_ajax_url': upload_ajax_url,
    })


def _update_daily_activity(user, qa_is_new=True):
    """Increment daily question count when a new Q&A row is created.

    Study minutes and streaks are updated in _update_user_analytics via
    _update_study_time_and_streak to avoid double-counting the same interaction.
    """
    if not qa_is_new:
        return
    today = date.today()
    activity, _ = DailyActivity.objects.get_or_create(
        user=user,
        date=today,
        defaults={'questions_asked': 0}
    )
    activity.questions_asked += 1
    activity.save()


def _update_user_analytics(user, qa_was_created=True):
    """Update user analytics after a Q&A interaction."""
    from testa_app.study_assistant_views import _update_study_time_and_streak

    analytics, _ = UserAnalytics.objects.get_or_create(user=user)

    # Keep total in sync with the database (covers creates, updates, deletes elsewhere)
    analytics.total_questions = QuestionAnswer.objects.filter(user=user).count()

    analytics.last_active = timezone.now()

    # Keep favorite_course in sync with most-asked course
    top_course = (
        QuestionAnswer.objects
        .filter(user=user)
        .exclude(course='')
        .values('course')
        .annotate(n=Count('id'))
        .order_by('-n')
        .first()
    )
    if top_course:
        analytics.favorite_course = top_course['course']

    analytics.save()

    # New Q&A only: one minute toward streak / totals (matches dashboard study charts)
    if qa_was_created:
        try:
            _update_study_time_and_streak(user, 60)
        except Exception:
            pass



@login_required
def all_questions(request):
    all_qa_list = QuestionAnswer.objects.filter(user=request.user).order_by('-created_at')

    return render(request, 'testa_app/all_questions.html', {
        'all_qa_list': all_qa_list,
    })


@require_http_methods(["POST", "DELETE"])
@login_required
def delete_question_answer(request, qa_id):
    """Delete a single Q&A entry. User must own it. Returns JSON."""
    try:
        qa = QuestionAnswer.objects.get(id=qa_id, user=request.user)
        qa.delete()
        try:
            _update_user_analytics(request.user, qa_was_created=False)
        except Exception:
            pass
        return JsonResponse({'success': True, 'deleted_id': qa_id})
    except QuestionAnswer.DoesNotExist:
        return JsonResponse({'success': False, 'error': 'Not found'}, status=404)
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@require_http_methods(["POST", "DELETE"])
@login_required
def delete_all_question_answers(request):
    """Delete all Q&A history for the current user. Returns JSON."""
    try:
        deleted_count, _ = QuestionAnswer.objects.filter(user=request.user).delete()
        try:
            _update_user_analytics(request.user, qa_was_created=False)
        except Exception:
            pass
        return JsonResponse({'success': True, 'deleted_count': deleted_count})
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


from docx import Document
from pptx import Presentation
import os

def get_file_text(file):
    extension = file.name.split('.')[-1].lower()
    
    if extension == 'pdf':
        return extract_text_from_pdf(file)
    elif extension == 'docx':
        return extract_text_from_docx(file)
    elif extension == 'pptx':
        return extract_text_from_pptx(file)
    elif extension == 'txt':
        return extract_text_from_txt(file)
    elif extension == 'ppt':
        return extract_text_from_txt(file)
    else:
        raise ValueError('Unsupported file type')

def extract_text_from_pdf(file):
    text = ""
    pdf_reader = PdfReader(file)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_txt(file):
    return file.read().decode('utf-8', errors='replace')



@login_required
def pdf_upload(request):
    if request.method == 'POST':
        form = PDFUploadForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_doc = form.save(commit=False)
            pdf_doc.uploaded_by = request.user
            if not pdf_doc.title:
                pdf_doc.title = request.FILES['file'].name
            pdf_doc.save()
            
            file = request.FILES['file']
            raw_text = get_file_text(file)
            text_chunks = get_text_chunks(raw_text)
            get_vector_store(text_chunks)
            messages.success(request, f'Successfully uploaded "{pdf_doc.title}"! You can now ask questions about this document.')
            return redirect('pdf_upload')
    else:
        form = PDFUploadForm()
    documents = PDFDocument.objects.filter(uploaded_by=request.user).order_by('-uploaded_at')
    return render(request, 'testa_app/pdf_upload.html', {'form': form, 'documents': documents})


@login_required
def upload_document_ajax(request):
    """Upload a document from the question/chat page (AJAX). Returns JSON."""
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Method not allowed'}, status=405)
    file = request.FILES.get('file')
    if not file:
        return JsonResponse({'success': False, 'error': 'No file provided'}, status=400)
    allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
    ext = file.name.split('.')[-1].lower()
    if ext not in allowed_extensions:
        return JsonResponse({'success': False, 'error': 'Unsupported file type. Use PDF, DOCX, PPTX, or TXT.'}, status=400)
    try:
        raw_text = get_file_text(file)
        text_chunks = get_text_chunks(raw_text)
        pdf_doc = PDFDocument(
            file=file,
            title=file.name,
            uploaded_by=request.user
        )
        pdf_doc.save()
        get_vector_store(text_chunks)
        return JsonResponse({
            'success': True,
            'title': pdf_doc.title or file.name,
            'message': f'"{pdf_doc.title or file.name}" is ready. Ask anything about it below.'
        })
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)





def register(request):
    from django.contrib import messages

    from .email_service import send_welcome_email
    from .models import UserProfile

    if request.method == 'POST':
        form = UserRegisterForm(request.POST)
        if form.is_valid():
            user = form.save()
            UserProfile.objects.create(user=user, email_verified=False)
            try:
                send_welcome_email(user, request)
            except Exception:
                messages.warning(
                    request,
                    'Account created, but we could not send the welcome email. '
                    'Use “Resend verification” on the login page.',
                )
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password1')
            user = authenticate(username=username, password=password)
            login(request, user)
            messages.success(
                request,
                'Welcome to Testa StudyBuddy! Check your email to verify your account.',
            )
            return redirect('question_answer')
    else:
        form = UserRegisterForm()
    return render(request, 'testa_app/register.html', {'form': form})




def login_view(request):
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password')
            user = authenticate(username=username, password=password)
            if user is not None:
                login(request, user)
                return redirect('question_answer')
    else:
        form = AuthenticationForm()
    return render(request, 'testa_app/login.html', {'form': form})



from django.http import JsonResponse
from django.views.decorators.http import require_POST
from .models import QuestionAnswer, Vote
import json

@require_http_methods(["PATCH"])
@login_required
def vote(request):
    try:
        data = json.loads(request.body)
        qa_id = data.get('qa_id')
        action = data.get('action')
        question_answer = QuestionAnswer.objects.get(id=qa_id)
        vote, created = Vote.objects.get_or_create(
            user=request.user, 
            question_answer=question_answer,
            defaults={'vote_type': action}
        )

        if not created: 
            if vote.vote_type == action:
                return JsonResponse({'error': 'You have already voted this way'}, status=400)
            else:
                vote.vote_type = action
                vote.save()

               
                if action == 'upvote':
                    question_answer.upvotes += 1
                    question_answer.downvotes -= 1
                else:
                    question_answer.upvotes -= 1
                    question_answer.downvotes += 1

        else:
            if action == 'upvote':
                question_answer.upvotes += 1
            elif action == 'downvote':
                question_answer.downvotes += 1

        question_answer.save()

        return JsonResponse({
            'upvotes': question_answer.upvotes,
            'downvotes': question_answer.downvotes
        })

    except QuestionAnswer.DoesNotExist:
        return JsonResponse({'error': 'QuestionAnswer not found'}, status=404)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def about(request):
    return render(request, 'about.html')
from django.contrib.auth import logout




from django.contrib.auth import logout

def logout_view(request):
    logout(request)
    return redirect('login')



def health(request):
    """Lightweight health check for Render (no DB)."""
    return JsonResponse({'status': 'ok', 'service': 'testa-studybuddy'})


def index(request):
    return render(request, 'index.html')


@login_required
def recommendations(request):
    """Smart recommendations view"""
    from .recommendation_engine import RecommendationEngine
    
    # Generate recommendations
    engine = RecommendationEngine(request.user)
    recommendations_list = engine.generate_recommendations(limit=20)
    
    # Get daily focus (highest priority)
    daily_focus = engine.get_daily_focus()
    
    # Filter by type if requested
    rec_type = request.GET.get('type', '')
    if rec_type:
        recommendations_list = [r for r in recommendations_list if r.recommendation_type == rec_type]
    
    return render(request, 'testa_app/recommendations.html', {
        'recommendations': recommendations_list,
        'daily_focus': daily_focus,
        'filter_type': rec_type,
    })


@login_required
def complete_recommendation(request, rec_id):
    """Mark recommendation as complete"""
    if request.method == 'POST':
        from .models import Recommendation
        recommendation = get_object_or_404(Recommendation, id=rec_id, user=request.user)
        recommendation.is_completed = True
        recommendation.completed_at = timezone.now()
        recommendation.save()
        return JsonResponse({'success': True})
    return JsonResponse({'success': False}, status=400)


@login_required
def profile(request):
    """User profile page"""
    user = request.user
    
    # Get or create user analytics
    analytics, created = UserAnalytics.objects.get_or_create(user=user)
    
    # Get user statistics
    total_questions = QuestionAnswer.objects.filter(user=user).count()
    total_bookmarks = Bookmark.objects.filter(user=user).count()
    total_quizzes = Quiz.objects.filter(created_by=user).count()
    total_flashcards = Flashcard.objects.filter(user=user).count()
    total_documents = PDFDocument.objects.filter(uploaded_by=user).count()
    
    # Get recent activity
    recent_questions = QuestionAnswer.objects.filter(user=user).order_by('-created_at')[:5]
    recent_bookmarks = Bookmark.objects.filter(user=user).order_by('-created_at')[:5]
    
    # Get topic mastery
    topic_mastery = TopicMastery.objects.filter(user=user).order_by('-mastery_level')[:5]
    
    # Get favorite course
    favorite_course = analytics.favorite_course or "Not set"
    
    # Calculate account age
    account_age_days = (timezone.now() - user.date_joined).days
    
    context = {
        'user': user,
        'analytics': analytics,
        'total_questions': total_questions,
        'total_bookmarks': total_bookmarks,
        'total_quizzes': total_quizzes,
        'total_flashcards': total_flashcards,
        'total_documents': total_documents,
        'recent_questions': recent_questions,
        'recent_bookmarks': recent_bookmarks,
        'topic_mastery': topic_mastery,
        'favorite_course': favorite_course,
        'account_age_days': account_age_days,
    }
    
    return render(request, 'testa_app/profile.html', context)

