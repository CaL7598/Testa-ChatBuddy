"""
Utility functions for Testa studyBuddy application

Provides utilities for document processing, AI interactions, and study tools
for students across all academic disciplines.
"""
import os
import json
from pathlib import Path
from datetime import datetime, timedelta
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import PromptTemplate
from .bytez_client import BytezClient, EmbeddingClient, get_bytez_client


def get_file_text(file):
    """Extract text from various file types (PDF, DOCX, PPTX, TXT)"""
    extension = file.name.split('.')[-1].lower()
    
    if extension == 'pdf':
        return extract_text_from_pdf(file)
    elif extension == 'docx':
        return extract_text_from_docx(file)
    elif extension == 'pptx':
        return extract_text_from_pptx(file)
    elif extension == 'txt':
        return extract_text_from_txt(file)
    else:
        raise ValueError('Unsupported file type')


def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    text = ""
    pdf_reader = PdfReader(file)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text


def extract_text_from_docx(file):
    """Extract text from DOCX file"""
    doc = Document(file)
    return '\n'.join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file):
    """Extract text from PPTX file"""
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_txt(file):
    """Extract text from TXT file"""
    return file.read().decode('utf-8', errors='replace')


def get_text_chunks(text, chunk_size=50000, chunk_overlap=1000):
    """Split text into chunks for vector database"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    chunks = text_splitter.split_text(text)
    return chunks


# Cap RAG context sent to the LLM (fewer tokens → faster responses, lower cost)
_RAG_CONTEXT_CHAR_LIMIT = int(os.environ.get("RAG_CONTEXT_CHAR_LIMIT", "14000"))


def _truncate_rag_context(text: str, limit: int = None) -> str:
    if not text:
        return text
    lim = limit if limit is not None else _RAG_CONTEXT_CHAR_LIMIT
    if len(text) <= lim:
        return text
    return text[:lim] + "\n\n[Retrieved context truncated for response speed.]"


def _get_faiss_index_path():
    """Return absolute path for FAISS index so save/load always use same location."""
    from django.conf import settings
    path = Path(settings.BASE_DIR) / "faiss_index"
    path.mkdir(parents=True, exist_ok=True)
    return str(path)


# --- Module-level cache: embedding model and vector store ---
_embedding_client_cache = None
_local_embeddings_cache = None
_vector_store_cache = None
_vector_store_mtime = None


def _get_local_embeddings():
    """Return a cached LocalEmbeddings instance, loading the model only once."""
    global _embedding_client_cache, _local_embeddings_cache
    if _local_embeddings_cache is not None:
        return _local_embeddings_cache

    from langchain_core.embeddings import Embeddings

    if _embedding_client_cache is None:
        _embedding_client_cache = EmbeddingClient()

    _client = _embedding_client_cache

    class LocalEmbeddings(Embeddings):
        def embed_documents(self, texts):
            return _client.embed_documents(texts)
        def embed_query(self, text):
            return _client.embed_text(text)

    _local_embeddings_cache = LocalEmbeddings()
    return _local_embeddings_cache


def load_vector_store(index_path=None):
    """Load FAISS vector store; uses an in-process cache to avoid reloading from disk."""
    global _vector_store_cache, _vector_store_mtime
    if index_path is None:
        index_path = _get_faiss_index_path()

    index_file = Path(index_path) / "index.faiss"
    if not index_file.exists():
        return None

    try:
        mtime = index_file.stat().st_mtime
    except OSError:
        return None

    # Return cached store when the file hasn't changed
    if _vector_store_cache is not None and _vector_store_mtime == mtime:
        return _vector_store_cache

    embeddings = _get_local_embeddings()
    try:
        vector_store = FAISS.load_local(
            index_path,
            embeddings,
            allow_dangerous_deserialization=True
        )
        _vector_store_cache = vector_store
        _vector_store_mtime = mtime
        return vector_store
    except Exception as e:
        print(f"Error loading vector store: {e}")
        return None


def get_vector_store(text_chunks, index_path=None):
    """Create or update FAISS vector store; merges new chunks into existing index."""
    global _vector_store_cache, _vector_store_mtime
    if index_path is None:
        index_path = _get_faiss_index_path()

    embeddings = _get_local_embeddings()
    existing = load_vector_store(index_path)

    if existing is not None and text_chunks:
        existing.add_texts(text_chunks)
        existing.save_local(index_path)
        # Update the in-memory cache and invalidate mtime so reload picks up changes
        _vector_store_cache = existing
        _vector_store_mtime = None
        return existing

    if text_chunks:
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local(index_path)
        _vector_store_cache = vector_store
        _vector_store_mtime = None
        return vector_store

    return existing


def get_conversational_chain():
    """Get conversational chain for question answering using Bytez"""
    client = get_bytez_client()

    def answer_with_context(context, question):
        """Answer question with context using Bytez"""
        context = _truncate_rag_context(context)
        return client.answer_question(question, context=context)
    
    return answer_with_context


class QuizGenerator:
    """Generate quizzes from documents using AI via Bytez"""
    
    def __init__(self, api_key=None):
        self.client = BytezClient(api_key=api_key) if api_key else get_bytez_client()

    def generate_quiz(self, document_text, topic, num_questions=5, difficulty="medium"):
        """Generate quiz questions from document text"""
        prompt = f"""Generate a quiz with {num_questions} questions on the topic of {topic} from this content:

{document_text}

Difficulty: {difficulty}

Return ONLY a JSON object with this structure:
{{
  "title": "Quiz title",
  "questions": [
    {{
      "question": "Question text",
      "type": "mcq",
      "options": ["A", "B", "C", "D"],
      "correct_answer": "A",
      "explanation": "Why A is correct"
    }}
  ]
}}

Ensure questions test understanding, not just memorization.
"""
        try:
            response = self.client.generate_text(
                prompt,
                system_prompt="You are an educational quiz generator. Always return valid JSON.",
                temperature=0.7,
                max_length=2048
            )
            content = response
            
            # Try to parse JSON
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            if json_start != -1 and json_end > json_start:
                json_str = content[json_start:json_end]
                return json.loads(json_str)
            return None
        except Exception as e:
            print(f"Error generating quiz: {e}")
            return None


class FlashcardGenerator:
    """Generate flashcards from documents using AI via Bytez"""

    def __init__(self, api_key=None):
        self.client = BytezClient(api_key=api_key) if api_key else get_bytez_client()
    
    def generate_flashcards(self, document_text, topic, num_cards=10):
        """Generate flashcards from document text"""
        prompt = f"""Create {num_cards} flashcards from this educational content about {topic}:

{document_text}

Return ONLY a JSON object:
{{
  "flashcards": [
    {{
      "front": "Question or term",
      "back": "Answer or definition"
    }}
  ]
}}

Focus on key concepts, definitions, and important facts.
"""
        try:
            response = self.client.generate_text(
                prompt,
                system_prompt="You are an educational flashcard generator. Always return valid JSON.",
                temperature=0.7,
                max_length=2048
            )
            content = response
            
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            if json_start != -1 and json_end > json_start:
                json_str = content[json_start:json_end]
                return json.loads(json_str)
            return None
        except Exception as e:
            print(f"Error generating flashcards: {e}")
            return None


class SummaryGenerator:
    """Generate summaries from documents using Bytez"""

    def __init__(self, api_key=None):
        self.client = BytezClient(api_key=api_key) if api_key else get_bytez_client()
    
    def generate_summary(self, document_text, summary_type="concise"):
        """Generate summary based on type: concise, detailed, or bullet_points"""
        type_instructions = {
            "concise": "Provide 2-3 paragraphs covering main points.",
            "detailed": "Cover all topics and subtopics comprehensively.",
            "bullet_points": "Create organized bullet list of key concepts."
        }
        
        prompt = f"""Summarize this educational content in {summary_type} format:

{document_text}

{type_instructions.get(summary_type, type_instructions['concise'])}
"""
        try:
            response = self.client.generate_text(
                prompt,
                system_prompt="You are an educational content summarizer.",
                temperature=0.3,
                max_length=2048
            )
            return response
        except Exception as e:
            print(f"Error generating summary: {e}")
            return None
    
    def generate_study_guide(self, document_text, topic):
        """Generate comprehensive study guide"""
        prompt = f"""Create a comprehensive study guide for the topic: {topic}

Content:
{document_text}

The study guide should include:
1. Overview of the topic
2. Key concepts list
3. Important formulas (if applicable)
4. Study tips
5. Common mistakes to avoid
6. Practice questions (3-5)

Format it in clear sections with headings.
"""
        try:
            response = self.client.generate_text(
                prompt,
                system_prompt="You are an educational study guide creator. Format responses with clear sections and headings.",
                temperature=0.3,
                max_length=4096
            )
            return response
        except Exception as e:
            print(f"Error generating study guide: {e}")
            return None
